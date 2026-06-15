#!/usr/bin/env python3
# Generates the GShare introduction deck, written from the product's point of view: a GPU sharing
# platform on Kubernetes. The slide content is Korean, matching docs/gshare.html and the published
# PDF and PowerPoint.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SHOT = "/home/jn/gShare/docs/screenshots"  # 콘솔 데모 스크린샷

# Palette, matching the landing page in docs/gshare.html.
BG    = RGBColor(0x0B,0x10,0x20)
BG2   = RGBColor(0x12,0x19,0x33)
CARD  = RGBColor(0x16,0x1E,0x3D)
LINE  = RGBColor(0x26,0x30,0x5A)
BRAND = RGBColor(0x6C,0x8C,0xFF)
TEAL  = RGBColor(0x37,0xE0,0xC8)
ACCENT= RGBColor(0xFF,0xB8,0x6B)
RED   = RGBColor(0xFF,0x7A,0x9C)
WHITE = RGBColor(0xFF,0xFF,0xFF)
FG    = RGBColor(0xEE,0xF1,0xFB)
MUTED = RGBColor(0x9A,0xA6,0xCF)
DARK  = RGBColor(0x22,0x22,0x22)
FONT  = "Malgun Gothic"

prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=None):
    s = prs.slides.add_slide(BLANK)
    if bg is not None:
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,SW,SH)
        r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def box(s, x,y,w,h, fill, line=None, lw=1.25, rounded=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x,y,w,h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    return shp

def tb(s, x,y,w,h, lines, size=18, color=FG, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, space=6):
    bx=s.shapes.add_textbox(x,y,w,h); tf=bx.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    if isinstance(lines,str): lines=[lines]
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        text, kw = ln if isinstance(ln,tuple) else (ln,{})
        r=p.add_run(); r.text=text
        f=r.font; f.size=Pt(kw.get("size",size)); f.bold=kw.get("bold",bold)
        f.color.rgb=kw.get("color",color); f.name=FONT
    return bx

def header(s, title, kicker, accent=TEAL):
    box(s, 0,0, SW, SH, BG)                 # 본문 슬라이드는 다크 배경 통일
    box(s, 0,0, Inches(0.16), SH, accent)   # 좌측 액센트 바
    tb(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
       [(kicker.upper(),{"size":13,"color":accent,"bold":True})])
    tb(s, Inches(0.7), Inches(0.85), Inches(12.1), Inches(0.95),
       [(title,{"size":30,"color":WHITE,"bold":True})])
    tb(s, Inches(0.7), Inches(7.06), Inches(11), Inches(0.35),
       [("GShare · 쿠버네티스 GPU 공유 플랫폼",{"size":10,"color":MUTED})])

def shot(s, fname, x, y, w, label, sub):
    # Place a 1440x900 console screenshot inside a browser frame, with a caption underneath.
    imgw = w - Inches(0.12)
    imgh = Emu(int(imgw * 900 / 1440))
    barh = Inches(0.30)
    fh = barh + imgh + Inches(0.1)
    box(s, x, y, w, fh, CARD, LINE, 1.25, True)        # 프레임
    box(s, x+Inches(0.03), y+Inches(0.03), w-Inches(0.06), barh, BG2)  # 타이틀 바
    for i,c in enumerate([RGBColor(0xFF,0x5F,0x57),RGBColor(0xFE,0xBC,0x2E),RGBColor(0x28,0xC8,0x40)]):
        d=s.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.18)+Inches(0.2)*i, y+Inches(0.11), Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb=c; d.line.fill.background(); d.shadow.inherit=False
    s.shapes.add_picture(os.path.join(SHOT,fname), x+Inches(0.06), y+barh, width=imgw, height=imgh)
    tb(s, x+Inches(0.02), y+fh+Inches(0.1), w, Inches(0.35), [(label,{"size":16,"color":WHITE,"bold":True})])
    tb(s, x+Inches(0.02), y+fh+Inches(0.47), w, Inches(0.5), [(sub,{"size":12.5,"color":MUTED})])

# ───────────────────────── 1. Title ─────────────────────────
s=slide(BG)
box(s, 0,0, SW, Inches(0.16), BRAND)
tb(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.5),
   [("KUBERNETES GPU SHARING PLATFORM",{"size":15,"color":TEAL,"bold":True})])
tb(s, Inches(0.9), Inches(1.7), Inches(11.7), Inches(2.0),
   [("GShare",{"size":62,"color":WHITE,"bold":True}),
    ("하나의 GPU 클러스터를 조직 전체가 나눠 씁니다",{"size":28,"color":BRAND,"bold":True})])
tb(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(1.1),
   [("여러 조직·부서·사용자가 하나의 GPU 클러스터를 분할(fractional)·전용(exclusive)으로 나눠 쓰고,",{"size":18,"color":RGBColor(0xCF,0xD6,0xF5)}),
    ("크레딧으로 과금·관리하는 GPU 공유 플랫폼.",{"size":18,"color":RGBColor(0xCF,0xD6,0xF5)})])
chips=["분할·전용 혼용","크레딧 기반 과금","유휴 GPU 자동 회수","멀티 클러스터 연동"]
x=Inches(0.9)
for c in chips:
    units=sum(1.0 if ('가'<=ch<='힣') else 0.55 for ch in c)  # 한글은 넓게, 영문·공백··는 좁게
    w=Inches(0.6+0.215*units)
    box(s, x, Inches(5.5), w, Inches(0.55), CARD, LINE, rounded=True)
    t=tb(s, x, Inches(5.5), w, Inches(0.55), [(c,{"size":15,"color":TEAL,"bold":True})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t.text_frame.word_wrap=False  # 한 줄 고정
    x += w + Inches(0.2)
tb(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.4),
   [("GPU를 직접 쪼개는 대신, 그 위에서 누가·얼마나·어떤 비용으로 쓰는지를 책임지는 제어 영역.",
     {"size":15,"color":MUTED})])

# ───────────────────────── 2. Problem versus approach ─────────────────────────
s=slide(); header(s,"GPU는 비싼데, 늘 누군가는 놀고 있습니다","왜 필요한가",ACCENT)
# Problem panel
box(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.6), RGBColor(0x2A,0x16,0x24), RGBColor(0x5A,0x27,0x42), 1.5, True)
tb(s, Inches(1.0), Inches(2.2), Inches(5.3), Inches(0.5), [("기존 방식의 한계",{"size":20,"color":RED,"bold":True})])
prob=["1인 1카드 전용 할당 → 점유 시간 대부분 GPU 유휴",
      "단순 time-slicing·분할은 비용 귀속·한도 관리가 없음",
      "유휴 세션이 카드를 붙잡아 대기열만 길어짐",
      "부서·조직별 예산 통제·감사 추적이 어려움",
      "큰 작업이 들어갈 빈 카드가 조각나 사라짐(단편화)"]
tb(s, Inches(1.0), Inches(2.75), Inches(5.3), Inches(3.7),
   [("✕  "+p,{"size":16,"color":RGBColor(0xDD,0xE3,0xF7)}) for p in prob], space=11)
# Approach panel
box(s, Inches(6.85), Inches(2.0), Inches(5.8), Inches(4.6), RGBColor(0x13,0x26,0x24), RGBColor(0x22,0x56,0x4F), 1.5, True)
tb(s, Inches(7.15), Inches(2.2), Inches(5.3), Inches(0.5), [("GShare의 접근",{"size":20,"color":TEAL,"bold":True})])
sol=["분할·전용을 한 클러스터에서 혼용",
     "점유율(occupancy) 기반 과금 — 쓴 만큼 크레딧 차감",
     "유휴 GPU는 자동 일시정지로 즉시 회수, 작업은 보존",
     "계층 크레딧·예산 게이트·자원 정책으로 한도 관리",
     "점유율을 고려한 bin-packing으로 빈 카드 보존"]
tb(s, Inches(7.15), Inches(2.75), Inches(5.3), Inches(3.7),
   [("✓  "+p,{"size":16,"color":RGBColor(0xDD,0xE3,0xF7)}) for p in sol], space=11)

# ───────────────────────── 3. Value ─────────────────────────
s=slide(); header(s,"한눈에 보는 GShare","핵심 가치",TEAL)
stats=[("2가지 모드","분할 vGPU · 전용 풀카드",BRAND),
       ("4계층 크레딧","시스템 → 조직 → 부서 → 개인",TEAL),
       ("유휴 자동 회수","idle GPU 자동 일시정지",ACCENT),
       ("멀티 클러스터","외부 GPU 클러스터 연동",RED)]
w=Inches(2.85); gap=Inches(0.25); x=Inches(0.7)
for big,sub,col in stats:
    box(s, x, Inches(2.2), w, Inches(2.3), CARD, LINE, rounded=True)
    box(s, x, Inches(2.2), w, Inches(0.12), col)
    tb(s, x+Inches(0.05), Inches(2.7), w-Inches(0.1), Inches(1.0),
       [(big,{"size":24,"color":col,"bold":True})], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x+Inches(0.1), Inches(3.7), w-Inches(0.2), Inches(0.7),
       [(sub,{"size":15,"color":MUTED})], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    x += w + gap
box(s, Inches(0.7), Inches(5.0), Inches(11.92), Inches(1.5), BG2, LINE, rounded=True)
tb(s, Inches(1.1), Inches(5.2), Inches(11.2), Inches(1.2),
   [("실제 GPU 분할은 HAMi(분할 vGPU)와 전용 풀카드가 맡고,",{"size":19,"color":FG}),
    ("GShare는 그 위에서 배치·과금·회수·권한을 하나의 콘솔로 통합 운영합니다.",{"size":19,"color":WHITE,"bold":True})],
   anchor=MSO_ANCHOR.MIDDLE)

# ───────────────────────── 4. Features ─────────────────────────
s=slide(); header(s,"분할 위에 얹은 운영 지능","핵심 기능",BRAND)
feats=[("점유율을 고려한 bin-packing","VRAM·코어 두 축의 잔여를 함께 보는 best-fit 배치. 빈 카드를 보존해 전용·대형 작업 수용력 유지.",BRAND),
       ("유휴 자동 일시정지","유휴(util≈0) 세션을 종료가 아닌 자동 일시정지로 회수. 세션·볼륨·크레딧은 보존, 언제든 재개.",TEAL),
       ("무손실 일시정지","cuda-checkpoint로 GPU 메모리 상태를 보존한 채 카드를 반환하고 재개 시 그대로 복원. 모델 재로딩 비용 0.",ACCENT),
       ("계층 크레딧 · 예산","시스템→조직→부서→개인 배분. 세션은 본인 지갑에서만 과금, 범위별 예산 게이트·월 리필.",BRAND),
       ("자원 정책 · 한도","동시 실행 수·자원 총량·유휴 타임아웃을 정책으로 통제. 사용자→부서→조직→공통 순으로 해석.",TEAL),
       ("대화형 세션","VSCode·JupyterLab·웹터미널을 1회용 링크로 즉시 접속. 볼륨·스냅샷·감사 로그까지 콘솔 한 곳에서.",ACCENT)]
cw=Inches(3.84); ch=Inches(2.15); gx=Inches(0.2); gy=Inches(0.25)
x0=Inches(0.7); y0=Inches(2.0)
for i,(t,d,col) in enumerate(feats):
    cx=x0+(cw+gx)*(i%3); cy=y0+(ch+gy)*(i//3)
    box(s, cx, cy, cw, ch, CARD, LINE, rounded=True)
    box(s, cx, cy, Inches(0.1), ch, col)
    tb(s, cx+Inches(0.28), cy+Inches(0.2), cw-Inches(0.5), Inches(0.6),
       [(t,{"size":18,"color":WHITE,"bold":True})])
    tb(s, cx+Inches(0.28), cy+Inches(0.82), cw-Inches(0.5), Inches(1.25),
       [(d,{"size":14,"color":MUTED})])

# ───────────────────────── 5. How it works ─────────────────────────
s=slide(); header(s,"세션 생성 한 번에 일어나는 일","동작 방식",TEAL)
tb(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
   [("요청은 허용 검사를 순서대로 통과해야 기동되고, 그때부터 점유율 기반 과금이 시작됩니다.",{"size":17,"color":MUTED})])
steps=[("STEP 1","선택","오퍼링(GPU 모델·티어)·이미지·볼륨을 고르면 예상 크레딧을 표시."),
       ("STEP 2","허용 검사","정책 한도 → 예산 → 크레딧 hold → 점유율을 고려한 VRAM 예약 순으로 통과."),
       ("STEP 3","기동","제어 영역이 CR 적용 → operator가 파드·서비스·인그레스 생성."),
       ("STEP 4","과금 · 회수","running 콜백부터 과금 시작, 유휴 시 자동 일시정지로 GPU 반환.")]
cw=Inches(2.85); x=Inches(0.7); y=Inches(2.7)
for i,(n,t,d) in enumerate(steps):
    box(s, x, y, cw, Inches(3.0), CARD, LINE, rounded=True)
    tb(s, x+Inches(0.25), y+Inches(0.25), cw-Inches(0.5), Inches(0.4),
       [(n,{"size":13,"color":BRAND,"bold":True})])
    tb(s, x+Inches(0.25), y+Inches(0.7), cw-Inches(0.5), Inches(0.6),
       [(t,{"size":21,"color":WHITE,"bold":True})])
    tb(s, x+Inches(0.25), y+Inches(1.45), cw-Inches(0.5), Inches(1.4),
       [(d,{"size":15,"color":MUTED})])
    if i<3:
        tb(s, x+cw-Inches(0.02), y+Inches(1.1), Inches(0.3), Inches(0.6),
           [("›",{"size":34,"color":BRAND,"bold":True})], align=PP_ALIGN.CENTER)
    x += cw + Inches(0.22)

# ───────────────────────── 6. Demo: the user console ─────────────────────────
s=slide(); header(s,"실제 콘솔 — 사용자 화면","데모 · 사용자 콘솔",TEAL)
tb(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5),
   [("세션을 직접 만들고, 접속하고, 크레딧을 확인하는 셀프서비스 콘솔.",{"size":17,"color":MUTED})])
ushots=[("02-user-seoyeon-dashboard.png","대시보드","내 자원·크레딧·GPU 가용성 한눈에"),
        ("04-user-seoyeon-session-new.png","새 세션 만들기","오퍼링·티어·이미지 선택 + 예상 크레딧"),
        ("06-user-seoyeon-session-connect.png","접속","VSCode·Jupyter·터미널 1회용 링크")]
w=Inches(3.84); x=Inches(0.7)
for fn,lb,sub in ushots:
    shot(s, fn, x, Inches(2.5), w, lb, sub); x += w + Inches(0.2)

# ───────────────────────── 7. Architecture ─────────────────────────
s=slide(); header(s,"제어 영역 · 실행 영역 분리","아키텍처",BRAND)
tb(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
   [("돈·상태 결정은 제어 영역이 책임지고, 실제 쿠버네티스 실행은 클러스터별 operator가 맡습니다.",{"size":17,"color":MUTED})])
layers=[("콘솔  (frontend · React SPA)","사용자·관리자 화면 (REST/SSE)",BRAND),
        ("제어 영역  (backend · FastAPI)","gshare-api: 인증·권한·허용검사·CRUD · gshare-worker: 과금·예산 롤업·대기열·리필",TEAL),
        ("operator  (Go · 클러스터별)","GShareSession CR → Pod·Service·Ingress 생성 (gshare-sessions ns)",ACCENT),
        ("HAMi + GPU 노드","실제 GPU 분할(gpumem/gpucores) · 전용 풀카드",RED)]
y=Inches(2.55); lw_=Inches(9.0)
for i,(t,d,col) in enumerate(layers):
    box(s, Inches(0.7), y, lw_, Inches(0.92), CARD, col, 1.5, True)
    box(s, Inches(0.7), y, Inches(0.12), Inches(0.92), col)
    tb(s, Inches(1.0), y+Inches(0.1), lw_-Inches(0.4), Inches(0.4),
       [(t,{"size":17,"color":WHITE,"bold":True})])
    tb(s, Inches(1.0), y+Inches(0.5), lw_-Inches(0.4), Inches(0.4),
       [(d,{"size":12.5,"color":MUTED})])
    if i<3:
        tb(s, Inches(0.7), y+Inches(0.9), lw_, Inches(0.22),
           [("▼",{"size":13,"color":col,"bold":True})], align=PP_ALIGN.CENTER)
    y += Inches(1.12)
# Side panel: the state stores
box(s, Inches(10.0), Inches(2.55), Inches(2.6), Inches(2.0), BG2, LINE, 1.25, True)
tb(s, Inches(10.2), Inches(2.7), Inches(2.3), Inches(1.8),
   [("상태 저장소",{"size":14,"color":TEAL,"bold":True}),
    ("postgres",{"size":15,"color":WHITE,"bold":True}),
    ("상태 · 원장(ledger)",{"size":12,"color":MUTED}),
    ("redis",{"size":15,"color":WHITE,"bold":True}),
    ("큐 · 토큰 · 락",{"size":12,"color":MUTED})], space=3)
box(s, Inches(10.0), Inches(4.75), Inches(2.6), Inches(1.55), BG2, LINE, 1.25, True)
tb(s, Inches(10.2), Inches(4.9), Inches(2.3), Inches(1.3),
   [("멀티 클러스터",{"size":14,"color":ACCENT,"bold":True}),
    ("제어 영역이 kubeconfig로 외부 GPU 클러스터에도 세션을 적용",{"size":13,"color":MUTED})], space=3)

# ───────────────────────── 8. Roles ─────────────────────────
s=slide(); header(s,"전사부터 개인까지, 범위에 맞게","역할 기반 운영",TEAL)
tb(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
   [("하나의 콘솔에서 역할 범위에 따라 메뉴와 데이터가 제한됩니다.",{"size":17,"color":MUTED})])
rows=[("역할","범위","할 수 있는 일",True,TEAL),
      ("슈퍼관리자","전사","조직·오퍼링·정책·클러스터·노드·이미지 관리, 크레딧 충전",False,BRAND),
      ("조직관리자","자기 조직","부서·사용자 관리, 조직→부서 예산 배분, 조직 세션 관제·감사",False,TEAL),
      ("부서관리자","자기 부서","부서원 관리, 부서→개인 배분, 부서 세션 관제",False,ACCENT),
      ("일반 사용자","본인","세션 생성·접속, 지갑·크레딧 요청, 볼륨·스냅샷",False,RED)]
y=Inches(2.45); rh=Inches(0.86)
for i,(a,b,c,hdr,col) in enumerate(rows):
    if hdr:
        tb(s, Inches(0.9), y, Inches(2.3), rh, [(a,{"size":13,"color":TEAL,"bold":True})], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, Inches(3.3), y, Inches(2.0), rh, [(b,{"size":13,"color":TEAL,"bold":True})], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, Inches(5.4), y, Inches(7.0), rh, [(c,{"size":13,"color":TEAL,"bold":True})], anchor=MSO_ANCHOR.MIDDLE)
        box(s, Inches(0.7), y+rh-Inches(0.02), Inches(11.9), Pt(1.5), LINE)
        y += Inches(0.55)
        continue
    box(s, Inches(0.7), y, Inches(0.1), rh-Inches(0.12), col)
    tb(s, Inches(0.95), y, Inches(2.3), rh, [(a,{"size":17,"color":WHITE,"bold":True})], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(3.3), y, Inches(2.0), rh, [(b,{"size":15,"color":MUTED})], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(5.4), y, Inches(7.0), rh, [(c,{"size":15,"color":RGBColor(0xDD,0xE3,0xF7)})], anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(0.7), y+rh, Inches(11.9), Pt(1), LINE)
    y += rh

# ───────────────────────── 9. Demo: the admin console ─────────────────────────
s=slide(); header(s,"실제 콘솔 — 관리자 화면","데모 · 관리자 콘솔",BRAND)
tb(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5),
   [("클러스터 가동률·세션·크레딧을 하나의 콘솔에서 관제·배분.",{"size":17,"color":MUTED})])
ashots=[("16-superadmin-admin-dashboard.png","관리자 대시보드","가동률·VRAM 적재율·헬스 경보"),
        ("35-superadmin-admin-monitor.png","세션 관제","전체 세션·대기열 실시간(LIVE SSE)"),
        ("34-superadmin-admin-allocations.png","크레딧 배분","시스템→조직→부서 배분·월 리필")]
w=Inches(3.84); x=Inches(0.7)
for fn,lb,sub in ashots:
    shot(s, fn, x, Inches(2.5), w, lb, sub); x += w + Inches(0.2)

# ───────────────────────── 10. Credits and billing ─────────────────────────
s=slide(); header(s,"계층 크레딧으로 비용을 통제","크레딧 · 과금",ACCENT)
tb(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
   [("크레딧은 위에서 아래로 배분되고, 과금은 점유한 만큼 본인 지갑에서 차감됩니다.",{"size":17,"color":MUTED})])
tiers=[("시스템","슈퍼관리자가 충전",BRAND),
       ("조직","조직 예산 한도",TEAL),
       ("부서","부서 예산 한도",ACCENT),
       ("개인","세션이 과금되는 지갑",RED)]
cw=Inches(2.7); x=Inches(0.7); y=Inches(2.7)
for i,(t,d,col) in enumerate(tiers):
    box(s, x, y, cw, Inches(1.6), CARD, col, 1.5, True)
    tb(s, x, y+Inches(0.3), cw, Inches(0.6), [(t,{"size":22,"color":col,"bold":True})], align=PP_ALIGN.CENTER)
    tb(s, x, y+Inches(0.95), cw, Inches(0.5), [(d,{"size":13,"color":MUTED})], align=PP_ALIGN.CENTER)
    if i<3:
        tb(s, x+cw-Inches(0.04), y+Inches(0.45), Inches(0.34), Inches(0.6),
           [("›",{"size":32,"color":col,"bold":True})], align=PP_ALIGN.CENTER)
    x += cw + Inches(0.22)
pts=[("점유율 기반 과금","running 콜백부터 세션이 점유한 자원·시간만큼 크레딧 차감 — 쓴 만큼만 비용."),
     ("예산 게이트","세션 생성 전 범위별 예산을 검사해 초과를 차단·경보, 월 자동 리필 지원."),
     ("크레딧 hold","기동 직전 예상 비용을 선점유(hold)해 과다 사용·중복 기동을 방지.")]
y=Inches(4.7)
for t,d in pts:
    box(s, Inches(0.7), y, Inches(0.1), Inches(0.7), ACCENT)
    tb(s, Inches(0.95), y, Inches(3.0), Inches(0.7), [(t,{"size":17,"color":WHITE,"bold":True})], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(4.0), y, Inches(8.6), Inches(0.7), [(d,{"size":15,"color":RGBColor(0xDD,0xE3,0xF7)})], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.8)

# ───────────────────────── 11. Adoption ─────────────────────────
s=slide(); header(s,"5분 데모부터 멀티 클러스터까지","도입",BRAND)
opts=[("5분 데모","Docker Compose","제어 영역과 콘솔을 띄워 화면·흐름을 바로 체험. GPU 없이도 시연 가능.",BRAND),
      ("올인원 배포","GPU 클러스터","전체 스택(콘솔·제어 영역·operator·HAMi)을 GPU 클러스터에 배포해 실제 세션을 운영.",TEAL),
      ("멀티 클러스터","kubeconfig 연동","하나의 제어 영역에서 외부 GPU 클러스터들에 세션을 적용·과금·관제.",ACCENT)]
cw=Inches(3.84); x=Inches(0.7); y=Inches(2.3)
for t,h,d,col in opts:
    box(s, x, y, cw, Inches(3.4), CARD, LINE, rounded=True)
    box(s, x, y, cw, Inches(0.14), col)
    tb(s, x+Inches(0.3), y+Inches(0.45), cw-Inches(0.6), Inches(0.5), [(t,{"size":22,"color":col,"bold":True})])
    tb(s, x+Inches(0.3), y+Inches(1.05), cw-Inches(0.6), Inches(0.5), [(h,{"size":16,"color":WHITE,"bold":True})])
    box(s, x+Inches(0.3), y+Inches(1.65), cw-Inches(0.6), Pt(1), LINE)
    tb(s, x+Inches(0.3), y+Inches(1.85), cw-Inches(0.6), Inches(1.4), [(d,{"size":15,"color":MUTED})])
    x += cw + Inches(0.2)
box(s, Inches(0.7), Inches(6.0), Inches(11.92), Inches(0.75), BG2, LINE, rounded=True)
tb(s, Inches(0.7), Inches(6.0), Inches(11.92), Inches(0.75),
   [("시작하기 가이드 · 사용자/관리자 매뉴얼 — github.com/boanlab/gShare",{"size":15,"color":TEAL,"bold":True})],
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ───────────────────────── 12. Closing ─────────────────────────
s=slide(BG)
box(s, 0, Inches(3.25), SW, Inches(0.06), TEAL)
tb(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2),
   [("하나의 GPU 클러스터를, 조직 전체가, 통제 가능하게.",{"size":32,"color":WHITE,"bold":True})])
tb(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.2),
   [("GShare — 분할·전용을 한 클러스터에서, 크레딧으로 과금하고, 유휴는 자동 회수.",{"size":20,"color":BRAND})])
tb(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.6),
   [("github.com/boanlab/gShare · Apache-2.0",{"size":16,"color":MUTED})])

out="/home/jn/gShare/docs/gshare.pptx"
prs.save(out)
print("saved", out, "slides=", len(prs.slides._sldIdLst))
